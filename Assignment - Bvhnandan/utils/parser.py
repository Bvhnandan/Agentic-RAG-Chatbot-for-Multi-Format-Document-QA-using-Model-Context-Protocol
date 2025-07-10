import pandas as pd
import pptx
import docx
import markdown
from pathlib import Path
import fitz 

def parse_txt(file):
    return file.read().decode("utf-8")

def parse_md(file):
    return markdown.markdown(file.read().decode("utf-8"))

def parse_csv(file):
    df = pd.read_csv(file)
    return df.to_string(index=False)

def parse_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def parse_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(file):
    presentation = pptx.Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_file(file):
    ext = Path(file.name).suffix.lower()
    if ext == ".txt":
        return parse_txt(file)
    elif ext == ".md":
        return parse_md(file)
    elif ext == ".csv":
        return parse_csv(file)
    elif ext == ".pdf":
        return parse_pdf(file)
    elif ext == ".docx":
        return parse_docx(file)
    elif ext == ".pptx":
        return parse_pptx(file)
    else:
        return "Unsupported file format"
